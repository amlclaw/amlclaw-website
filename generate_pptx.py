#!/usr/bin/env python3
"""Generate AMLClaw pitch deck — polished version matching HTML."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme ──
BG = RGBColor(15, 15, 15)
CARD_BG = RGBColor(22, 22, 22)
CARD_BORDER = RGBColor(50, 50, 50)
FG = RGBColor(227, 227, 227)
ACCENT = RGBColor(138, 180, 248)
ACCENT_DIM = RGBColor(40, 60, 90)
DIM = RGBColor(130, 130, 130)
RED_M = RGBColor(248, 113, 113)
GREEN_M = RGBColor(110, 231, 183)
FONT = 'Arial'
MONO = 'Courier New'
W, H = Inches(13.333), Inches(7.5)
BASE = '/Users/max/Desktop/max-skills/amlclaw-website/'

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def add_slide():
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return sl


def txbox(sl, left, top, width, height, text, size=18, bold=False, color=FG,
          align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tf


def card(sl, x, y, w, h, fill=CARD_BG, border=CARD_BORDER):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = border
    shp.line.width = Pt(1)
    # Reduce corner rounding
    shp.adjustments[0] = 0.05
    return shp


def accent_card(sl, x, y, w, h):
    return card(sl, x, y, w, h, fill=RGBColor(18, 25, 40), border=RGBColor(60, 90, 140))


def divider_line(sl, x, y, w):
    """Thin horizontal accent line."""
    shp = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = ACCENT
    shp.line.fill.background()
    return shp


def slide_title(sl, text, subtitle=None):
    txbox(sl, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
          text, 36, True, FG, PP_ALIGN.LEFT)
    if subtitle:
        txbox(sl, Inches(1), Inches(1.1), Inches(11), Inches(0.5),
              subtitle, 18, False, DIM, PP_ALIGN.LEFT)
    divider_line(sl, Inches(1), Inches(1.35), Inches(2))


# ═══════════════════════════════════════════════════
# P1 — COVER
# ═══════════════════════════════════════════════════
sl = add_slide()
# Subtle accent line at top
divider_line(sl, Inches(4.5), Inches(1.5), Inches(4.3))
txbox(sl, Inches(1), Inches(1.8), Inches(11.3), Inches(1.8),
      "Humans Work for AI.\nAI Works for Compliance.",
      48, True, FG, PP_ALIGN.CENTER)
divider_line(sl, Inches(4.5), Inches(3.8), Inches(4.3))
txbox(sl, Inches(1), Inches(4.2), Inches(11.3), Inches(0.7),
      "AMLClaw — Open Source AI Compliance Platform",
      22, False, ACCENT, PP_ALIGN.CENTER)
# Bottom
txbox(sl, Inches(1), Inches(6.5), Inches(11.3), Inches(0.5),
      "github.com/amlclaw", 16, False, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P2 — PAIN POINTS
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "You're Still Paying Humans to Read PDFs?")

stats = [
    ("💸", "$500K–$1M / year", "Cost of a compliance team"),
    ("🐌", "4 Weeks", "To write one policy for one jurisdiction"),
    ("⏰", "Half a Day", "To screen a single address manually"),
    ("🙈", "80%", "Of crypto startups skip compliance entirely"),
]
for i, (emoji, big, sub) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    c = card(sl, x, y, Inches(2.8), Inches(3.8))
    txbox(sl, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.6),
          emoji, 36, False, FG, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.2), y + Inches(1.1), Inches(2.4), Inches(1.0),
          big, 24, True, ACCENT, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.2), y + Inches(2.4), Inches(2.4), Inches(1.2),
          sub, 14, False, DIM, PP_ALIGN.CENTER)

txbox(sl, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
      '"Most compliance teams cost more than the seed round."',
      16, True, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P3 — ASSEMBLY LINE
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "The Compliance Assembly Line", "Every step is manual. Every step costs money.")

roles = [
    ("👤 Lawyer", "Read regulations", "2 weeks"),
    ("👤 Compliance\n    Officer", "Write policy", "2 weeks"),
    ("👤 Engineer", "Build rules", "1 week"),
    ("👤 Analyst", "Screen addresses", "half day each"),
    ("👤 Manager", "Monitor", "spot checks"),
]
for i, (role, task, time) in enumerate(roles):
    x = Inches(0.5 + i * 2.55)
    c = card(sl, x, Inches(2.2), Inches(2.2), Inches(3.2))
    txbox(sl, x + Inches(0.15), Inches(2.4), Inches(1.9), Inches(0.9),
          role, 16, True, ACCENT, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.15), Inches(3.4), Inches(1.9), Inches(0.6),
          task, 14, False, FG, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.15), Inches(4.2), Inches(1.9), Inches(0.5),
          time, 13, False, RED_M, PP_ALIGN.CENTER)
    if i < 4:
        txbox(sl, x + Inches(2.15), Inches(3.2), Inches(0.5), Inches(0.5),
              "→", 24, True, DIM, PP_ALIGN.CENTER)

txbox(sl, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
      "5 humans  ×  12 months  =  💰💰💰",
      22, True, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P4 — FLIP IT (table style matching HTML)
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "What If We Flip It?")

tasks_list = ["Read regulations", "Write policies", "Build rules",
              "Screen addresses", "Monitor 24/7"]

# Table-style layout
tbl_shape = sl.shapes.add_table(6, 4, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.5))
tbl = tbl_shape.table
# Column widths
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.8)
tbl.columns[2].width = Inches(0.7)
tbl.columns[3].width = Inches(3.8)

# Header row
headers = ["", "OLD", "", "NEW"]
for c, h in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = FONT
        p.font.color.rgb = DIM if c == 1 else (ACCENT if c == 3 else DIM)
        p.alignment = PP_ALIGN.CENTER
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(25, 25, 25)

# Data rows
for r, task in enumerate(tasks_list):
    row_idx = r + 1
    cells = [
        (task, DIM, PP_ALIGN.RIGHT),
        ("👤 Humans", DIM, PP_ALIGN.CENTER),
        ("→", DIM, PP_ALIGN.CENTER),
        ("🤖 AI", ACCENT, PP_ALIGN.CENTER),
    ]
    for c, (txt, clr, al) in enumerate(cells):
        cell = tbl.cell(row_idx, c)
        cell.text = txt
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.bold = (c == 3)
            p.font.name = FONT
            p.font.color.rgb = clr
            p.alignment = al
        cell.fill.solid()
        if c == 3:
            cell.fill.fore_color.rgb = RGBColor(18, 28, 50)
        else:
            cell.fill.fore_color.rgb = RGBColor(20, 20, 20)

# Taglines
txbox(sl, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
      "Humans work FOR AI — train, guide, supervise",
      20, True, ACCENT, PP_ALIGN.CENTER)
txbox(sl, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
      "AI works FOR compliance — read, write, execute, monitor",
      20, True, ACCENT, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P5 — PIPELINE
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "One Pipeline. Full Compliance.")

steps = [
    ("📄", "Documents", "40+ regs\npre-loaded", "Human uploads"),
    ("📋", "Policies", "AI generates\nin 2 min", "Human reviews"),
    ("⚙️", "Rules", "AI converts\nto JSON", "Human approves"),
    ("🔍", "Screening", "On-chain trace\n+ rule matching", "Human decides"),
    ("📡", "Monitoring", "7×24\nautomated", "Human sleeps 😴"),
]
for i, (icon, name, desc, human) in enumerate(steps):
    x = Inches(0.5 + i * 2.55)
    c = accent_card(sl, x, Inches(1.8), Inches(2.2), Inches(3.8))
    txbox(sl, x + Inches(0.15), Inches(2.0), Inches(1.9), Inches(0.6),
          icon, 32, False, FG, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.15), Inches(2.6), Inches(1.9), Inches(0.5),
          name, 18, True, ACCENT, PP_ALIGN.CENTER)
    txbox(sl, x + Inches(0.15), Inches(3.3), Inches(1.9), Inches(0.9),
          desc, 13, False, FG, PP_ALIGN.CENTER)
    # Human role at bottom with green
    txbox(sl, x + Inches(0.15), Inches(4.5), Inches(1.9), Inches(0.6),
          human, 12, True, GREEN_M, PP_ALIGN.CENTER)
    if i < 4:
        txbox(sl, x + Inches(2.15), Inches(3.2), Inches(0.5), Inches(0.5),
              "→", 24, True, ACCENT, PP_ALIGN.CENTER)

txbox(sl, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
      "AI does 95% of the work. Humans review at key checkpoints.",
      16, False, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P6 — DEMO SCREENSHOTS (one per slide, full-size)
# ═══════════════════════════════════════════════════
imgs = ["documents.png", "policies.png", "rules-1.png",
        "screening-1.png", "screening-graph.png", "dashboard.png"]
labels = ["📄 Step 1: Load Regulations & Documents",
          "📋 Step 2: AI Generates Compliance Policies",
          "⚙️ Step 3: Convert Policies to Executable Rules",
          "🔍 Step 4: Screen Addresses Against Rules",
          "🕸️ Step 5: Visualize On-Chain Relationships",
          "📡 Step 6: Monitor Everything From One Dashboard"]
for i, (img, lbl) in enumerate(zip(imgs, labels)):
    sl = add_slide()
    # Title bar
    txbox(sl, Inches(0.5), Inches(0.3), Inches(8), Inches(0.5),
          "From Regulation to Screening in 5 Minutes", 20, True, DIM, PP_ALIGN.LEFT)
    # Step label
    txbox(sl, Inches(0.5), Inches(0.85), Inches(12), Inches(0.5),
          lbl, 24, True, ACCENT, PP_ALIGN.LEFT)
    # Progress dots
    for j in range(len(imgs)):
        dot = sl.shapes.add_shape(MSO_SHAPE.OVAL,
              Inches(11.5 + j * 0.3), Inches(0.45), Inches(0.15), Inches(0.15))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT if j == i else RGBColor(50, 50, 50)
        dot.line.fill.background()
    # Screenshot - large centered
    img_path = os.path.join(BASE, img)
    if os.path.exists(img_path):
        card(sl, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.7))
        sl.shapes.add_picture(img_path, Inches(0.9), Inches(1.6), Inches(11.5), Inches(5.5))

# ═══════════════════════════════════════════════════
# P7 — BUILT DIFFERENT
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "Built Different")

features = [
    ("🤖", "Multi-AI", "Claude · DeepSeek · Gemini\nSwitch anytime, no vendor lock-in"),
    ("📁", "No Database", "File-based JSON storage\nDeploy anywhere, backup = cp"),
    ("🔓", "Open Source", "MIT License\nEvery rule visible, full transparency"),
    ("📝", "Audit Trail", "Append-only JSONL\nTamper-resistant, compliance-ready"),
    ("🌍", "Bilingual", "English + 中文\nDark & Light theme"),
    ("🐳", "Docker Ready", "One command to deploy\nSelf-hosted, data stays yours"),
]
for i, (icon, title, desc) in enumerate(features):
    col, row = i % 3, i // 3
    x = Inches(0.7 + col * 4.1)
    y = Inches(1.7 + row * 2.7)
    c = card(sl, x, y, Inches(3.8), Inches(2.3))
    txbox(sl, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5),
          f"{icon}  {title}", 22, True, ACCENT, PP_ALIGN.LEFT)
    txbox(sl, x + Inches(0.3), y + Inches(0.9), Inches(3.2), Inches(1.2),
          desc, 14, False, DIM, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════
# P8 — TWO PRODUCTS
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "Two Products. Both Free. Both Yours.")

products = [
    ("🦅  AMLClaw Skill", "For AI Agents",
     "• Install in one sentence\n• Claude Code / OpenClaw / any framework\n• CLI-based, lightweight\n• Ships as SKILL.md entry point",
     "clawhub install amlclaw",
     "github.com/amlclaw/amlclaw"),
    ("🖥️  AMLClaw Dashboard", "For Humans",
     "• Full Web UI (Next.js)\n• Visual flow graphs & reports\n• Monitoring dashboard\n• Self-hostable, Docker-ready",
     "docker compose up -d",
     "github.com/amlclaw/amlclaw.com"),
]
for i, (title, audience, bullets, cmd, repo) in enumerate(products):
    x = Inches(0.7 + i * 6.3)
    c = accent_card(sl, x, Inches(1.7), Inches(5.8), Inches(5.0))
    txbox(sl, x + Inches(0.4), Inches(1.9), Inches(5.0), Inches(0.5),
          title, 26, True, ACCENT, PP_ALIGN.LEFT)
    txbox(sl, x + Inches(0.4), Inches(2.5), Inches(5.0), Inches(0.4),
          audience, 18, False, DIM, PP_ALIGN.LEFT)
    divider_line(sl, x + Inches(0.4), Inches(3.0), Inches(5.0))
    txbox(sl, x + Inches(0.4), Inches(3.2), Inches(5.0), Inches(2.0),
          bullets, 15, False, FG, PP_ALIGN.LEFT)
    # Command box
    cmd_card = card(sl, x + Inches(0.4), Inches(5.3), Inches(5.0), Inches(0.6),
                    fill=RGBColor(30, 35, 30), border=RGBColor(60, 90, 60))
    txbox(sl, x + Inches(0.6), Inches(5.35), Inches(4.6), Inches(0.5),
          f"$ {cmd}", 15, True, GREEN_M, PP_ALIGN.LEFT, font=MONO)
    txbox(sl, x + Inches(0.4), Inches(6.1), Inches(5.0), Inches(0.3),
          repo, 12, False, DIM, PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════
# P9 — THE MATH
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "The Math Is Embarrassingly Simple")

rows_data = [
    ("", "Traditional Team", "AMLClaw"),
    ("Annual Cost", "$500K – $1M", "~$50–200/mo AI API *"),
    ("Setup Time", "3–6 months", "5 minutes"),
    ("Policy Generation", "2–4 weeks", "2 minutes"),
    ("Address Screening", "Half a day", "< 5 minutes"),
    ("Monitoring", "Spot checks", "24/7 automated"),
    ("Transparency", "Black-box scores", "Open source, every rule visible"),
]

tbl_shape = sl.shapes.add_table(len(rows_data), 3,
                                 Inches(1.2), Inches(1.8), Inches(10.8), Inches(4.0))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(4.3)

for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15 if r > 0 else 16)
            p.font.name = FONT
            p.font.bold = (r == 0 or c == 0)
            p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            if r == 0:
                p.font.color.rgb = ACCENT
            elif c == 1:
                p.font.color.rgb = DIM
            elif c == 2:
                p.font.color.rgb = ACCENT
            else:
                p.font.color.rgb = FG
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = RGBColor(25, 30, 45)
        elif c == 2:
            cell.fill.fore_color.rgb = RGBColor(15, 22, 40)
        elif r % 2 == 0:
            cell.fill.fore_color.rgb = RGBColor(22, 22, 22)
        else:
            cell.fill.fore_color.rgb = RGBColor(18, 18, 18)

txbox(sl, Inches(1.2), Inches(6.2), Inches(10.8), Inches(0.4),
      "* AI API cost depends on usage. ~$50/mo covers most small-to-mid operations. Platform itself is free.",
      12, False, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P10 — ROADMAP
# ═══════════════════════════════════════════════════
sl = add_slide()
slide_title(sl, "What's Next")

# Now section
now_card = accent_card(sl, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.4))
txbox(sl, Inches(1.2), Inches(1.9), Inches(2), Inches(0.5),
      "NOW", 18, True, GREEN_M, PP_ALIGN.LEFT)
txbox(sl, Inches(1.2), Inches(2.4), Inches(10.8), Inches(0.5),
      "✅ 3 Jurisdictions (SG, HK, Dubai)     ✅ 40+ Regulations     ✅ Tron Chain",
      16, False, FG, PP_ALIGN.LEFT)

# Next section
next_card = card(sl, Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.8))
txbox(sl, Inches(1.2), Inches(3.6), Inches(2), Inches(0.5),
      "NEXT", 18, True, ACCENT, PP_ALIGN.LEFT)
items = "🔗 Bitcoin + Multi-chain Stablecoins\n🪪 KYC + Travel Rule  —  TrustIn acquired a KYC company, building identity + on-chain compliance\n🇪🇺 MiCA (EU)     🇺🇸 FinCEN (US)     🔌 Plugin System"
txbox(sl, Inches(1.2), Inches(4.1), Inches(10.8), Inches(1.0),
      items, 15, False, FG, PP_ALIGN.LEFT)

# Quote box
quote_card = card(sl, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.4),
                  fill=RGBColor(20, 20, 20), border=RGBColor(45, 45, 45))
txbox(sl, Inches(1.8), Inches(5.7), Inches(9.7), Inches(1.2),
      "Why Open Source?\nTrustIn already runs the SaaS. AMLClaw is the open-source edition —\nmaking AML accessible to everyone. Compliance shouldn't be a privilege.\nIt should be infrastructure.",
      14, True, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P11 — TRY IT NOW
# ═══════════════════════════════════════════════════
sl = add_slide()
txbox(sl, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8),
      "Try It Now. On Your Phone.  ⚡", 42, True, FG, PP_ALIGN.CENTER)
divider_line(sl, Inches(5), Inches(1.7), Inches(3.3))

steps_data = [
    ("1", "Open OpenClaw on your phone  🦞", FG),
    ("2", 'Tell your agent:\n"Search GitHub for amlclaw and install it"', FG),
    ("3", "Done. Your AI agent is now a compliance expert.", ACCENT),
]
for i, (num, text, color) in enumerate(steps_data):
    y = Inches(2.3 + i * 1.5)
    # Step number circle
    c = card(sl, Inches(2.5), y, Inches(0.6), Inches(0.6),
             fill=ACCENT_DIM, border=ACCENT)
    txbox(sl, Inches(2.5), y + Inches(0.05), Inches(0.6), Inches(0.5),
          num, 22, True, ACCENT, PP_ALIGN.CENTER)
    # Step text
    txbox(sl, Inches(3.5), y + Inches(0.05), Inches(7.5), Inches(1.0),
          text, 20, False, color, PP_ALIGN.LEFT)

txbox(sl, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
      '"You just deployed a compliance team in 10 seconds."',
      18, True, DIM, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# P12 — CLOSING
# ═══════════════════════════════════════════════════
sl = add_slide()
divider_line(sl, Inches(4.5), Inches(2.2), Inches(4.3))
txbox(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
      "The End Game of Compliance\nIsn't More People —\nIt's a Better System.",
      42, True, FG, PP_ALIGN.CENTER)
divider_line(sl, Inches(4.5), Inches(5.0), Inches(4.3))
txbox(sl, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
      "github.com/amlclaw", 20, True, ACCENT, PP_ALIGN.CENTER)
txbox(sl, Inches(1), Inches(6.1), Inches(11.3), Inches(0.4),
      "MIT License  ·  Free Forever", 16, False, DIM, PP_ALIGN.CENTER)


# ── Save ──
out = os.path.join(BASE, 'pitch.pptx')
prs.save(out)
print(f"✅ Saved: {out}")
